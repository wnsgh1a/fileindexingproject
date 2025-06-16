import os
import shutil
import hashlib
import difflib
from datetime import datetime
from collections import defaultdict, deque
from docx import Document
import pdfplumber
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup
import json
import subprocess

CANDIDATE_DIR = r"C:\Users\wnsgh\Desktop\삭제후보"
GUBOJEON_DIR = os.path.join(CANDIDATE_DIR, "구버전")
DUPLICATE_DIR = os.path.join(CANDIDATE_DIR, "중복파일")

def calculate_file_hash(file_path):
    sha256 = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            sha256.update(f.read())
        return sha256.hexdigest()
    except Exception as e:
        print(f"Error hashing {file_path}: {e}")
        return None

def move_to_category(file_path, category, reason=""):
    category_dir = os.path.join(CANDIDATE_DIR, category)
    os.makedirs(category_dir, exist_ok=True)
    destination = os.path.join(category_dir, os.path.basename(file_path))
    base, ext = os.path.splitext(destination)
    counter = 1
    while os.path.exists(destination):
        destination = f"{base}_{counter}{ext}"
        counter += 1
    shutil.move(file_path, destination)
    print(f"✅ {os.path.basename(file_path)} → {category_dir} 이유: {reason}")

def simplify_filename(filename):
    import re
    name, _ = os.path.splitext(filename.lower())
    name = re.sub(r"(ver|v)?[\._\-]?[0-9]+(\.[0-9]+)?", "", name)
    keywords = ["rev", "판본", "draft", "수정본", "최종", "복사본"]
    for keyword in keywords:
        name = name.replace(keyword, "")
    return name.strip().replace("_", "").replace(" ", "")

def convert_with_libreoffice(path):
    try:
        output_dir = os.path.dirname(path)
        subprocess.run([
            'soffice', '--headless', '--convert-to', 'docx', '--outdir', output_dir, path
        ], check=True)
        new_path = os.path.splitext(path)[0] + '.docx'
        if os.path.exists(new_path):
            return extract_text(new_path)
    except:
        return ""
    return ""

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext in [".txt", ".md", ".csv"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif ext == ".pdf":
            text = ""
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text
        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text += str(cell.value) + "\n"
            return text
        elif ext == ".pptx":
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif ext == ".html":
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "html.parser")
                return soup.get_text()
        elif ext == ".json":
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        elif ext in [".doc", ".ppt", ".xls"]:
            return convert_with_libreoffice(path)
    except:
        return ""
    return ""

def is_content_similar(file1, file2, threshold=0.85):
    try:
        text1 = extract_text(file1)
        text2 = extract_text(file2)
        similarity = difflib.SequenceMatcher(None, text1, text2).ratio()
        return similarity >= threshold
    except:
        return False

def build_similarity_clusters(similarity_groups):
    visited = set()
    clusters = []
    for file in similarity_groups:
        if file in visited:
            continue
        cluster = set()
        queue = deque([file])
        while queue:
            current = queue.popleft()
            if current in visited:
                continue
            visited.add(current)
            cluster.add(current)
            queue.extend(similarity_groups[current] - visited)
        if cluster:
            clusters.append(cluster)
    return clusters

def isolate_all(directory):
    print(f"\n📌 전체 폴더 기반 중복 및 구버전 정리 시작: {directory}\n")
    file_paths = []
    file_hashes = defaultdict(list)
    file_groups = defaultdict(list)
    duplicate_hashes = set()
    for root, _, files in os.walk(directory):
        for file in files:
            path = os.path.join(root, file)
            if not os.path.isfile(path):
                continue
            file_paths.append(path)
            file_hash = calculate_file_hash(path)
            if file_hash:
                file_hashes[file_hash].append(path)
    for hash_value, hash_group in file_hashes.items():
        if len(hash_group) <= 1:
            continue
        duplicate_hashes.update(hash_group)
        latest = max(hash_group, key=lambda x: os.path.getmtime(x))
        for f in hash_group:
            if f != latest:
                move_to_category(f, "중복파일", reason="전체 검사 기반 중복파일 정리")
    for path in file_paths:
        if not os.path.isfile(path) or path in duplicate_hashes:
            continue
        simplified = simplify_filename(os.path.basename(path))
        file_groups[simplified].append(path)
    for group, files in file_groups.items():
        if len(files) <= 1:
            continue
        similarity_graph = defaultdict(set)
        for i in range(len(files)):
            for j in range(i + 1, len(files)):
                if is_content_similar(files[i], files[j]):
                    similarity_graph[files[i]].add(files[j])
                    similarity_graph[files[j]].add(files[i])
        all_related = set()
        for k, v in similarity_graph.items():
            all_related.add(k)
            all_related.update(v)
        unclustered = set(files) - all_related
        for f in unclustered:
            similarity_graph[f] = set()
        clusters = build_similarity_clusters(similarity_graph)
        for cluster in clusters:
            latest = max(cluster, key=lambda x: os.path.getmtime(x))
            for f in cluster:
                if f != latest:
                    move_to_category(f, "구버전", reason="내용 유사 기반 구버전 정리")
    print("\n✅ 전체 정리가 완료되었습니다.\n")
